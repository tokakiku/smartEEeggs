# ==========================================
# 灵犀智课 - PPT 物理渲染引擎 (Banana JSON 绝对坐标架构)
# 特性：引入标准化中间层 | 纯白板绝对定位 | 像素(px)换算 | 完美图文混排 | 互动游戏超链接植入
# ==========================================

import os
import uuid
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from dotenv import load_dotenv

load_dotenv()
logger = logging.getLogger("PPT_Engine")

OUTPUT_DIR = "downloads/exports"

# 基础画布尺寸 (16:9 标准比例，映射为 1280x720 像素)
CANVAS_WIDTH_PX = 1280
CANVAS_HEIGHT_PX = 720
CANVAS_WIDTH_INCHES = 13.333
CANVAS_HEIGHT_INCHES = 7.5


def px_to_inches(px: int) -> float:
    """核心算法：将前端的像素坐标换算为 PPT 的英寸物理坐标"""
    return (px / CANVAS_WIDTH_PX) * CANVAS_WIDTH_INCHES


def hex_to_rgb(hex_color: str) -> RGBColor:
    """颜色转换器：将 #2c3e50 转换为底层 RGB 对象"""
    hex_color = hex_color.lstrip('#')
    if len(hex_color) == 6:
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
    return RGBColor(0, 0, 0)


# ==========================================
# 🌟 中间件：转换为 Banana JSON Schema
# ==========================================
def convert_to_banana_json(page_queue: list, project_id: int) -> dict:
    """
    将一维的图文队列，转化为带有精准排版坐标的标准化 Banana JSON
    """
    banana_json = {
        "metadata": {
            "title": f"项目 {project_id} 演示文稿",
            "author": "灵犀智课",
            "size": {"width": CANVAS_WIDTH_PX, "height": CANVAS_HEIGHT_PX},
        },
        "slides": []
    }

    for i, page in enumerate(page_queue):
        title = page.get("title", "未命名页面")
        content = page.get("content", "")
        image_path = page.get("image_path", "")
        game_url = page.get("game_url", "")  # 🌟 提取游戏链接

        slide_data = {
            "id": f"slide-{uuid.uuid4().hex[:8]}",
            "elements": []
        }

        # 1. 标题元素 (固定在左上方)
        slide_data["elements"].append({
            "type": "heading",
            "content": title,
            "style": {
                "fontSize": 40,
                "fontWeight": "bold",
                "color": "#1A202C",  # 深灰偏黑
                "position": {"x": 80, "y": 60},
                "width": 1120,
                "height": 80
            }
        })

        # 2. 判断是否有图片，动态计算正文宽度
        has_img = bool(image_path and os.path.exists(image_path))
        text_width = 600 if has_img else 1120

        # 3. 正文元素
        slide_data["elements"].append({
            "type": "text",
            "content": content,
            "style": {
                "fontSize": 22,
                "color": "#4A5568",  # 中度灰，适合阅读
                "width": text_width,
                "height": 450,  # 稍微缩短一点高度，给下方的游戏按钮留空间
                "position": {"x": 80, "y": 160}
            }
        })

        # 4. 图片元素 (如果有，放在右侧区域)
        if has_img:
            slide_data["elements"].append({
                "type": "image",
                "src": image_path,
                "style": {
                    "width": 450,
                    "height": 450,
                    "position": {"x": 750, "y": 160}
                }
            })

        # 5. 🌟 游戏超链接元素 (如果有，固定在左下角)
        if game_url:
            slide_data["elements"].append({
                "type": "game_link",
                "content": "🎮 点击启动课堂互动小游戏",
                "url": game_url,
                "style": {
                    "fontSize": 18,
                    "fontWeight": "bold",
                    "color": "#0070C0",  # 经典的超链接蓝色
                    "width": 400,
                    "height": 50,
                    "position": {"x": 80, "y": 620}  # 坐标定在左下方
                }
            })

        # 6. 页码元素 (右下角)
        slide_data["elements"].append({
            "type": "text",
            "content": f"{i + 1} / {len(page_queue)}",
            "style": {
                "fontSize": 14,
                "color": "#A0AEC0",
                "width": 100,
                "height": 30,
                "position": {"x": 1150, "y": 680}
            }
        })

        banana_json["slides"].append(slide_data)

    return banana_json


# ==========================================
# 🌟 物理渲染器：直接读取 JSON 坐标进行作画
# ==========================================
class BananaPPTGenerator:
    def __init__(self):
        # 创建一个全新的空白演示文稿（不再依赖 template.pptx）
        self.prs = Presentation()
        # 强制设置 16:9 画布比例
        self.prs.slide_width = Inches(CANVAS_WIDTH_INCHES)
        self.prs.slide_height = Inches(CANVAS_HEIGHT_INCHES)

    def render_element(self, slide, element: dict):
        """解析 JSON 里的 element，实现绝对定位渲染"""
        el_type = element.get("type")
        style = element.get("style", {})
        pos = style.get("position", {"x": 0, "y": 0})

        # 像素坐标转换为物理英寸
        left = Inches(px_to_inches(pos["x"]))
        top = Inches(px_to_inches(pos["y"]))
        width = Inches(px_to_inches(style.get("width", 800)))
        height = Inches(px_to_inches(style.get("height", 100)))

        # ----- 渲染文字类 (heading, text) -----
        if el_type in ["text", "heading"]:
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.text = element.get("content", "")

            # 排版与样式配置
            p.line_spacing = 1.3  # 加入呼吸感行距
            for run in p.runs:
                run.font.name = "Microsoft YaHei"
                run.font.size = Pt(style.get("fontSize", 20))
                run.font.bold = (style.get("fontWeight") == "bold")
                if "color" in style:
                    run.font.color.rgb = hex_to_rgb(style["color"])

        # ----- 🌟 渲染超链接类 (game_link) -----
        elif el_type == "game_link":
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()

            run.text = element.get("content", "")
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(style.get("fontSize", 18))
            run.font.bold = (style.get("fontWeight") == "bold")
            run.font.color.rgb = hex_to_rgb(style.get("color", "#0070C0"))
            run.font.underline = True  # 添加下划线

            game_url = element.get("url")
            if game_url:
                # 拼接完整服务器地址（本地联调时为 127.0.0.1:8000，上云后需替换为真实 IP/域名）
                full_game_link = f"http://127.0.0.1:8000{game_url}"
                run.hyperlink.address = full_game_link

        # ----- 渲染图片类 (image) -----
        elif el_type == "image":
            img_src = element.get("src", "")
            if os.path.exists(img_src):
                # 插入图片时可以自适应宽高比例
                slide.shapes.add_picture(img_src, left, top, width=width)

    def run(self, banana_json: dict, project_id: int) -> str:
        """核心执行：遍历 JSON 数组，逐页绘制"""
        try:
            slides_data = banana_json.get("slides", [])
            blank_layout = self.prs.slide_layouts[6]  # 使用系统的纯空白版式

            for slide_data in slides_data:
                # 每一页都是一张完全空白的画板
                slide = self.prs.slides.add_slide(blank_layout)

                # 遍历并绘制该页的所有元素
                elements = slide_data.get("elements", [])
                for el in elements:
                    self.render_element(slide, el)

            # 导出物理文件
            os.makedirs(OUTPUT_DIR, exist_ok=True)
            output_path = os.path.join(OUTPUT_DIR, f"presentation_{project_id}.pptx")
            self.prs.save(output_path)

            logger.info(f"✅ Banana 架构渲染完成（共 {len(slides_data)} 页）")
            return output_path
        except Exception as e:
            logger.error(f"❌ Banana 渲染崩溃: {str(e)}")
            raise e


# ==========================================
# 🌟 对外接口
# ==========================================
def generate_ppt_from_json(page_queue: list, project_id: int) -> str:
    """流水线包装：先转化成 Banana JSON，再进行物理渲染"""
    logger.info(f"🎨 [渲染引擎] 启动 Banana 架构，开始为项目 {project_id} 计算排版坐标...")
    banana_json = convert_to_banana_json(page_queue, project_id)
    engine = BananaPPTGenerator()
    return engine.run(banana_json, project_id)